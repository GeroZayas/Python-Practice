# These collections imports are necessary to avoid this error:
# AttributeError: module 'collections' has no attribute 'Container'
import collections
import collections.abc

# ----------------------------------------------------------------
from pptx import Presentation
from pptx.util import Inches
# ----------------------------------------------------------------

# Create a new PowerPoint presentation
prs = Presentation()

# Add a new slide
slide = prs.slides.add_slide(prs.slide_layouts[1])

# Define the table dimensions
rows = 3
cols = 4

# Define the table width and height
width = Inches(6.0)
height = Inches(4.0)

# Define the position of the table on the slide
left = top = Inches(1.0)

# Add a table to the slide
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Add data to the table
data = [['A', 'B', 'C', 'D'], [1, 2, 3, 4], [5, 6, 7, 8]]
for row in range(rows):
    for col in range(cols):
        table.cell(row, col).text = str(data[row][col])

# Save the PowerPoint presentation
prs.save('table_example.pptx')
