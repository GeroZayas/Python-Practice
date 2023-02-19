# These collections imports are necessary to avoid this error:
# AttributeError: module 'collections' has no attribute 'Container'
import collections
import collections.abc
# ----------------------------------------------------------------
from pptx import Presentation
# ----------------------------------------------------------------


prs = Presentation()
title_slide_layout = prs.slide_layouts[2]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, my name is Gero!"
subtitle.text = "this is pretty cool!"

for s in range(5):
    prs.slides.add_slide(prs.slide_layouts[5])

prs.save('Gero Zayas.pptx')
