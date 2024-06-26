# These collections imports are necessary to avoid this error:
# AttributeError: module 'collections' has no attribute 'Container'
# ----------------------------------------------------------------

import csv
from pptx import Presentation

# Create a new PowerPoint presentation
prs = Presentation()


# Add a title slide to the beginning of the presentation
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "GERO ZAYAS"
subtitle.text = "This is wonderful test"


# Open the CSV file and read the data
try:
    # Attempt to open the CSV file and read the data
    with open("questions_and_answers.csv", newline="") as f:
        reader = csv.reader(f)
        next(reader)  # Skip the header row
        for row in reader:
            # Create a new slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])

            # Add the question to the slide
            question = row[0]
            question_shape = slide.shapes.title
            question_shape.text = question

            # Add the answer to the slide
            answer = row[1]
            answer_shape = slide.placeholders[1]
            answer_shape.text = answer
except FileNotFoundError:
    print("The file 'questions_and_answers.csv' does not exist.")

# Save the PowerPoint presentation
prs.save("questions_and_answers.pptx")
